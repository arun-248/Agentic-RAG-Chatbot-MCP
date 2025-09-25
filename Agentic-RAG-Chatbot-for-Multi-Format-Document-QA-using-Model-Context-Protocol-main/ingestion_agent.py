# ingestion_agent.py
"""
Ingestion Agent

This module handles parsing and preprocessing of multi-format documents.
It converts raw files into clean text chunks that can be indexed for retrieval.

Supported formats:
- PDF (.pdf)
- PowerPoint (.pptx, .ppt)
- Word (.docx)
- CSV (.csv)
- Text / Markdown (.txt, .md)

Each chunk is returned with metadata and an MCP message is sent to the RetrievalAgent.
"""

import os
import logging
from typing import List, Dict
from mcp import send_mcp_message

import fitz  # PDF (PyMuPDF)
from pptx import Presentation
import pandas as pd
from docx import Document as DocxDocument

logging.basicConfig(level=logging.INFO)
CHUNK_SIZE = 800  # Developer note: adjust for better tradeoff between granularity and context length


# -------------------------------
# Helpers: Chunking & Parsing
# -------------------------------
def _chunk_text(text: str, size: int = CHUNK_SIZE) -> List[str]:
    """
    Split text into chunks of fixed size.
    """
    text = text.strip()
    if not text:
        return []
    return [text[i:i+size] for i in range(0, len(text), size)]


def parse_pdf(path: str) -> str:
    doc = fitz.open(path)
    pages = [p.get_text() for p in doc]
    return "\n".join(pages)


def parse_pptx(path: str) -> str:
    prs = Presentation(path)
    slides_text = []
    for idx, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        slides_text.append("\n".join(parts))
    return "\n\n".join(slides_text)


def parse_docx(path: str) -> str:
    doc = DocxDocument(path)
    paras = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n".join(paras)


def parse_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def parse_csv(path: str) -> str:
    df = pd.read_csv(path)
    # Developer note: limit to 50 rows for readability (prevents bloating chunks)
    return df.head(50).to_csv(index=False)


# -------------------------------
# Main: Process one document
# -------------------------------
def process_document(file_path: str, filename: str = None) -> List[Dict]:
    """
    Parse a file into text chunks + metadata.

    Returns:
        List[Dict]: List of {"text": ..., "meta": {"filename":..., "chunk_index":...}}
    Also sends MCP DOCUMENT_PARSED to RetrievalAgent.
    """
    if filename is None:
        filename = os.path.basename(file_path)

    ext = os.path.splitext(filename)[1].lower()
    try:
        if ext == ".pdf":
            text = parse_pdf(file_path)
        elif ext in [".pptx", ".ppt"]:
            text = parse_pptx(file_path)
        elif ext == ".docx":
            text = parse_docx(file_path)
        elif ext == ".csv":
            text = parse_csv(file_path)
        elif ext in [".txt", ".md"]:
            text = parse_txt(file_path)
        else:
            # Fallback: treat unknown format as text
            text = parse_txt(file_path)
    except Exception as e:
        logging.exception("Error parsing %s: %s", file_path, e)
        text = ""

    chunks = []
    for i, c in enumerate(_chunk_text(text)):
        chunks.append({"text": c, "meta": {"filename": filename, "chunk_index": i}})

    # Notify RetrievalAgent that parsing is complete
    send_mcp_message(
        sender="IngestionAgent",
        receiver="RetrievalAgent",
        type="DOCUMENT_PARSED",
        payload={"filename": filename, "n_chunks": len(chunks)}
    )

    logging.info("IngestionAgent: parsed %s -> %d chunks", filename, len(chunks))
    return chunks
